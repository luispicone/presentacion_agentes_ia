from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

OUT = Path(__file__).with_name('charla-excelencia-ia-organizacional.pptx')
ASSETS = Path(__file__).parent / 'assets'
LOGO = ASSETS / 'excelencia-logo.png'
FUNDECE = ASSETS / 'fundecelogo.jpg'
FPNC = ASSETS / 'fpnc.jpg'
IPACE = ASSETS / 'ipace.jpg'

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# EXCELENCIA-inspired palette: white + anthracite + colored diagonal accents from the X mark
WHITE = RGBColor(255, 255, 255)
ANTH = RGBColor(55, 55, 57)
GREY = RGBColor(105, 105, 108)
LIGHT = RGBColor(244, 246, 248)
LINE = RGBColor(225, 228, 232)
YELLOW = RGBColor(246, 187, 38)
ORANGE = RGBColor(239, 111, 40)
FUCHSIA = RGBColor(211, 47, 122)
TEAL = RGBColor(0, 173, 181)
DARKTEAL = RGBColor(0, 119, 130)
BLUEGREY = RGBColor(40, 58, 75)

COLORS = [YELLOW, ORANGE, FUCHSIA, TEAL]
FONT = 'Aptos'


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=24, color=ANTH, bold=False, align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, kicker=None):
    if kicker:
        add_textbox(slide, 0.72, 0.45, 8.5, 0.32, kicker.upper(), 10.5, TEAL, True)
    add_textbox(slide, 0.68, 0.75, 9.8, 0.78, title, 27, ANTH, True)
    add_logo(slide)
    add_accent_rule(slide)


def add_logo(slide):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.25), Inches(0.32), width=Inches(2.35))


def add_accent_rule(slide):
    x = 0.72
    y = 1.62
    for i, c in enumerate(COLORS):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + i * 0.42), Inches(y), Inches(0.32), Inches(0.055))
        shp.fill.solid(); shp.fill.fore_color.rgb = c
        shp.line.fill.background()


def add_footer(slide, num):
    add_textbox(slide, 0.72, 7.05, 4.4, 0.2, 'EXCELENCIA · IA organizacional', 8.5, GREY)
    add_textbox(slide, 12.15, 7.05, 0.45, 0.2, f'{num:02d}', 8.5, GREY, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.94), Inches(11.9), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_bullets(slide, x, y, w, h, items, size=17.5, color=ANTH, bullet_color=TEAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.line_spacing = 1.04
        # Use a colored dot as manual bullet for reliable rendering.
        p.text = '• ' + item
    return box


def card(slide, x, y, w, h, title, body, accent=TEAL, title_size=15.5, body_size=12.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_textbox(slide, x+0.25, y+0.18, w-0.45, 0.35, title, title_size, ANTH, True)
    add_textbox(slide, x+0.25, y+0.67, w-0.45, h-0.82, body, body_size, GREY)
    return shape


def quote(slide, x, y, w, h, text, accent=FUCHSIA):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT
    shape.line.fill.background()
    add_textbox(slide, x+0.28, y+0.22, w-0.56, h-0.42, text, 18, ANTH, True, line_spacing=1.05)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def diagonal_band(slide, x, y, w, h, color, transparency=5):
    shp = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.fill.transparency = transparency
    shp.line.fill.background()
    return shp


def pill(slide, x, y, w, text, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.48))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    add_textbox(slide, x+0.08, y+0.13, w-0.16, 0.22, text, 10.5, WHITE, True, align=PP_ALIGN.CENTER)

# 1
slide = blank(); set_bg(slide)
diagonal_band(slide, 8.9, -0.15, 0.75, 7.95, YELLOW, 8)
diagonal_band(slide, 9.65, -0.15, 0.75, 7.95, ORANGE, 8)
diagonal_band(slide, 10.4, -0.15, 0.75, 7.95, FUCHSIA, 8)
diagonal_band(slide, 11.15, -0.15, 0.75, 7.95, TEAL, 8)
if LOGO.exists(): slide.shapes.add_picture(str(LOGO), Inches(0.72), Inches(0.55), width=Inches(3.2))
add_textbox(slide, 0.78, 2.0, 7.85, 1.55, 'IA organizacional:\nde información confiable a agentes gobernados', 32, ANTH, True, line_spacing=0.92)
add_textbox(slide, 0.82, 3.75, 6.9, 0.65, 'Cómo preparar documentación, datos y procesos para incorporar IA con valor real.', 18.5, GREY)
pill(slide, 0.82, 5.0, 1.35, '30 min', TEAL); pill(slide, 2.32, 5.0, 2.5, 'CEOs y directores', FUCHSIA); pill(slide, 5.0, 5.0, 2.1, 'Gestión + IA', ORANGE)
add_textbox(slide, 0.82, 6.45, 4.7, 0.26, 'Luis Picone · EXCELENCIA', 12.5, GREY)
add_footer(slide, 1)

# 2
slide = blank(); set_bg(slide); add_title(slide, 'La tesis', 'Punto de partida')
quote(slide, 0.85, 2.15, 5.3, 1.55, 'La IA no empieza por el algoritmo ni por el agente. Empieza por la calidad de la información sobre la que va a trabajar.', TEAL)
card(slide, 6.65, 2.05, 2.0, 2.7, 'Documentación', 'Vigente, trazable y fiel a la realidad operativa.', YELLOW)
card(slide, 8.95, 2.05, 2.0, 2.7, 'Datos', 'Íntegros, consistentes, comparables y gobernados.', ORANGE)
card(slide, 11.25, 2.05, 1.45, 2.7, 'Gobierno', 'Fuentes, permisos, criterios y validación.', FUCHSIA)
add_bullets(slide, 0.95, 4.35, 5.7, 1.6, ['Sin base confiable, la IA amplifica inconsistencias.', 'Con base confiable, la IA puede acelerar análisis, trazabilidad y decisiones.'], 16.5)
add_footer(slide, 2)

# 3
slide = blank(); set_bg(slide); add_title(slide, 'Antes de IA: documentación fiel a la realidad', 'Base documental')
add_bullets(slide, 0.9, 2.05, 5.55, 3.4, ['La documentación debe representar cómo opera realmente la organización.', 'Los repositorios deben ser robustos, actualizados, accesibles y trazables.', 'Una IA puede responder con forma profesional aunque la fuente sea débil.', 'Si el proceso real y el documentado no coinciden, la IA trabaja sobre una ficción ordenada.'], 17)
quote(slide, 7.0, 2.15, 5.2, 1.25, 'La primera pregunta no es qué IA usar. Es si existe información confiable para que la IA pueda trabajar.', FUCHSIA)
card(slide, 7.0, 4.05, 5.2, 1.45, 'Riesgo ejecutivo', 'Decisiones apoyadas en documentos no vigentes, evidencia difícil de auditar y pérdida de confianza por problemas de base.', ORANGE, 15, 13)
add_footer(slide, 3)

# 4
slide = blank(); set_bg(slide); add_title(slide, 'Datos íntegros: la otra mitad de la base', 'Integridad de datos')
card(slide, 0.85, 2.05, 2.75, 2.1, 'Completos', 'Cobertura suficiente para analizar sin sesgos críticos.', YELLOW)
card(slide, 3.9, 2.05, 2.75, 2.1, 'Consistentes', 'Mismas reglas, definiciones y criterios entre áreas.', ORANGE)
card(slide, 6.95, 2.05, 2.75, 2.1, 'Trazables', 'Origen, responsable, actualización y evidencia.', FUCHSIA)
card(slide, 10.0, 2.05, 2.75, 2.1, 'Gobernados', 'Permisos, usos, controles y validación.', TEAL)
quote(slide, 1.2, 4.85, 10.9, 1.05, 'La IA puede acelerar análisis; no puede reemplazar la responsabilidad organizacional sobre la calidad de los datos.', TEAL)
add_footer(slide, 4)

# 5
slide = blank(); set_bg(slide); add_title(slide, 'El riesgo de incorporar IA sobre una base débil', 'Riesgo')
risks = ['Automatizar errores existentes', 'Dar velocidad a procesos mal definidos', 'Usar documentos no vigentes como fuente de verdad', 'Mezclar criterios entre áreas', 'Exponer información sensible sin permisos claros', 'Generar respuestas difíciles de validar']
for i, r in enumerate(risks):
    x = 0.85 + (i % 3) * 4.05
    y = 2.05 + (i // 3) * 1.35
    card(slide, x, y, 3.55, 0.95, f'{i+1}', r, COLORS[i % 4], title_size=14, body_size=12.5)
quote(slide, 1.05, 5.45, 11.1, 0.82, 'Sin base confiable, la IA no transforma la organización: amplifica sus inconsistencias.', FUCHSIA)
add_footer(slide, 5)

# 6
slide = blank(); set_bg(slide); add_title(slide, 'Dónde puede ayudar EXCELENCIA antes del agente', 'Rol de EXCELENCIA')
steps = [
    ('Diagnosticar', 'madurez documental y de datos'),
    ('Ordenar', 'repositorios, procesos y fuentes'),
    ('Gobernar', 'criterios, permisos, validación y riesgos'),
    ('Priorizar', 'casos de uso por valor y criticidad'),
]
for i, (t,b) in enumerate(steps):
    x = 0.9 + i * 3.1
    card(slide, x, 2.35, 2.55, 1.8, t, b, COLORS[i], 16, 13)
    if i < 3:
        add_textbox(slide, x+2.62, 2.95, 0.35, 0.25, '→', 22, GREY, True, align=PP_ALIGN.CENTER)
quote(slide, 1.05, 5.05, 11.1, 0.9, 'Antes de diseñar agentes, hay que preparar el terreno: documentación, datos, procesos, criterios y gobierno.', TEAL)
add_footer(slide, 6)

# 7
slide = blank(); set_bg(slide); add_title(slide, 'De usar IA a diseñar Agentes de IA', 'Madurez')
card(slide, 0.9, 2.05, 3.55, 2.65, 'Uso puntual de IA', 'Una persona pregunta, obtiene una respuesta y resuelve una tarea aislada.', YELLOW, 16, 13.2)
card(slide, 4.9, 2.05, 3.55, 2.65, 'IA asistida con buena base', 'La organización usa fuentes confiables y criterios comunes.', ORANGE, 16, 13.2)
card(slide, 8.9, 2.05, 3.55, 2.65, 'Agente de IA gobernado', 'Sistema diseñado para una tarea o proceso, con propósito, fuentes, reglas, herramientas, límites, trazabilidad, validación y métricas.', TEAL, 16, 12.3)
quote(slide, 1.25, 5.35, 10.45, 0.82, 'El agente no es el primer paso. Es el paso superador cuando la organización ya sabe qué información, proceso y criterio quiere gobernar.', FUCHSIA)
add_footer(slide, 7)

# 8
slide = blank(); set_bg(slide); add_title(slide, 'Procesos donde la IA puede agregar valor', 'Alcance potencial')
areas = ['Dirección y gobierno', 'Estrategia y planeamiento', 'Calidad y sistemas de gestión', 'Operaciones y seguridad', 'Comercial y clientes', 'Finanzas y control', 'Legales y compliance', 'Personas y cultura', 'Tecnología y sistemas', 'Compras y proveedores', 'Comunicación y sostenibilidad', 'Gestión documental']
for i, a in enumerate(areas):
    x = 0.85 + (i % 4) * 3.1
    y = 2.0 + (i // 4) * 0.82
    c = COLORS[i % 4]
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.68), Inches(0.55))
    shp.fill.solid(); shp.fill.fore_color.rgb = LIGHT
    shp.line.color.rgb = LINE
    add_textbox(slide, x+0.12, y+0.15, 2.44, 0.22, a, 10.5, ANTH, True, align=PP_ALIGN.CENTER)
quote(slide, 1.05, 5.2, 11.1, 0.92, 'La pregunta relevante no es si la IA puede aportar valor, sino qué información necesita, qué tarea conviene delegar, con qué autonomía y con qué controles.', TEAL)
add_footer(slide, 8)

# 9
slide = blank(); set_bg(slide); add_title(slide, 'Caso único · Documentación, calidad, ISO y PNC', 'Caso aplicado')
flow = [('1', 'Base documental\nanalizable'), ('2', 'Cruce contra\nISO y PNC'), ('3', 'Evidencias, brechas\ny oportunidades'), ('4', 'Validación\nexperta')]
for i,(n,t) in enumerate(flow):
    x = 0.95 + i*3.05
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.05), Inches(0.62), Inches(0.62))
    circ.fill.solid(); circ.fill.fore_color.rgb = COLORS[i]
    circ.line.fill.background()
    add_textbox(slide, x+0.14, 2.18, 0.35, 0.2, n, 14, WHITE, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.78, 2.0, 1.9, 0.7, t, 14.5, ANTH, True)
    if i < 3: add_textbox(slide, x+2.54, 2.22, 0.28, 0.2, '→', 18, GREY, True)
card(slide, 1.0, 3.45, 5.15, 1.35, 'Subejemplo ISO', 'Requisitos, evidencia, cobertura, brechas y acciones de cierre.', TEAL, 15.5, 12.8)
card(slide, 6.75, 3.45, 5.15, 1.35, 'Subejemplo PNC', 'Modelo de excelencia, fortalezas, oportunidades de mejora e insumos de evaluación.', FUCHSIA, 15.5, 12.8)
quote(slide, 1.05, 5.35, 11.1, 0.82, 'La IA no reemplaza el criterio experto; lo potencia cuando trabaja sobre documentación confiable.', FUCHSIA)
add_footer(slide, 9)

# 10
slide = blank(); set_bg(slide); add_title(slide, 'Subejemplo ISO · Matriz de gaps VMOS', 'Caso aplicado')
add_textbox(slide, 0.86, 1.78, 10.6, 0.28, 'Muestra real tomada de matriz-gaps-iso-vmos.xlsx · filas Excel 14 a 23 · ids 13 a 22 · cláusula 6', 11.5, GREY, False)
card(slide, 0.85, 2.18, 2.75, 1.28, 'Alcance', '10 requisitos ISO de planificación: riesgos, ambiente, SST, legales, objetivos y cambios.', TEAL, 15, 11.6)
card(slide, 3.85, 2.18, 2.75, 1.28, 'Resultado', '3 cubiertos y 7 parciales. La matriz distingue evidencia suficiente de brechas documentales.', ORANGE, 15, 11.6)
card(slide, 6.85, 2.18, 2.75, 1.28, 'Evidencias', 'PRO-EXC-RIE, PRO-EXC-AMB, PRO-EXC-SOP, PRO-EXC-REQ y proceso de Excelencia.', FUCHSIA, 15, 11.4)
card(slide, 9.85, 2.18, 2.75, 1.28, 'Brechas', 'Faltan registros/matrices integradas para aspectos, SST, objetivos SIG y cambios.', YELLOW, 15, 11.4)

rows = [
    ('Cubierto', 'Riesgos y oportunidades; requisitos legales; planificación de acciones', TEAL),
    ('Parcial', 'Aspectos ambientales; peligros/riesgos/oportunidades SST', ORANGE),
    ('Parcial', 'Objetivos del sistema y planificación para lograrlos', FUCHSIA),
    ('Parcial', 'Gestión del cambio: consecuencias, recursos, responsables e impactos SST', YELLOW),
]
y0 = 3.82
for i, (estado, detalle, color) in enumerate(rows):
    y = y0 + i * 0.44
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(y), Inches(1.28), Inches(0.31))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    add_textbox(slide, 1.02, y+0.07, 1.14, 0.13, estado, 8.0, WHITE, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.42, y+0.04, 4.05, 0.22, detalle, 9.2, ANTH, False)

card(slide, 6.75, 3.72, 5.65, 1.72, 'Valor del agente', 'Identifica documentos y evidencia; ubica secciones, citas o registros; propone cobertura o brecha; justifica el criterio; ordena acciones de cierre y conserva trazabilidad para validación experta.', TEAL, 14.5, 10.2)

quote(slide, 0.95, 5.95, 11.45, 0.62, 'Mensaje para la charla: el agente acelera lectura, cruce y trazabilidad; el experto conserva el criterio de cobertura y la decisión final.', TEAL)
add_footer(slide, 10)

# 11
slide = blank(); set_bg(slide); add_title(slide, 'Subejemplo PNC · Matriz de gaps VMOS', 'Caso aplicado')
add_textbox(slide, 0.86, 1.78, 10.9, 0.28, 'Muestra real tomada de matriz-gaps-pnc-vmos.xlsx · filas Excel 14 a 22 · ids 13 a 21 · Liderazgo / Factor 1.3', 11.0, GREY, False)
card(slide, 0.85, 2.18, 2.75, 1.28, 'Alcance', '9 aspectos del factor 1.3: planeamiento estratégico y operativo.', TEAL, 15, 11.6)
card(slide, 3.85, 2.18, 2.75, 1.28, 'Resultado', '9 cubiertos, 0 parciales, 0 no cubiertos según evidencia documental cargada.', ORANGE, 15, 11.2)
card(slide, 6.85, 2.18, 2.75, 1.28, 'Evidencias', 'Estrategia, Sustentabilidad, Finanzas, Personas, Riesgos, Comunicaciones e Innovación.', FUCHSIA, 15, 11.0)
card(slide, 9.85, 2.18, 2.75, 1.28, 'Salida útil', 'Cobertura por aspecto + fuente documental para revisión experta.', YELLOW, 15, 11.6)

rows = [
    ('Cubierto', 'a-b: horizonte, información decisoria y objetivos estratégicos'),
    ('Cubierto', 'c-d: innovación y despliegue de objetivos en planes operativos'),
    ('Cubierto', 'e-g: riesgos, recursos e indicadores clave de desempeño'),
    ('Cubierto', 'h-i: comunicación de estrategias y revisión de planes/objetivos'),
]
y0 = 3.82
for i, (estado, detalle) in enumerate(rows):
    y = y0 + i * 0.44
    color = TEAL
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(y), Inches(1.28), Inches(0.31))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    add_textbox(slide, 1.02, y+0.07, 1.14, 0.13, estado, 8.0, WHITE, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.42, y+0.04, 4.05, 0.22, detalle, 9.2, ANTH, False)

card(slide, 6.75, 3.72, 5.65, 1.72, 'Valor del agente', 'Lee el modelo; ordena evidencias por criterio, factor y aspecto; identifica cobertura documental, fuentes y trazabilidad; apoya diagnósticos y prepara insumos para revisión experta.', FUCHSIA, 14.5, 10.2)

quote(slide, 0.95, 5.95, 11.45, 0.62, 'Mensaje para la charla: el agente ordena la cobertura PNC por aspecto; el experto conserva el juicio de madurez.', FUCHSIA)
add_footer(slide, 11)

# 12
slide = blank(); set_bg(slide); add_title(slide, 'Cierre: la secuencia de madurez', 'Síntesis ejecutiva')
sequence = [('Información confiable', 'documentación vigente + datos íntegros'), ('IA con gobierno', 'fuentes, permisos, validación y trazabilidad'), ('Agentes gobernados', 'tareas/procesos con métricas y resultados')]
for i,(t,b) in enumerate(sequence):
    x = 1.0 + i*4.0
    card(slide, x, 2.15, 3.25, 1.55, t, b, COLORS[i], 15.5, 12.6)
    if i < 2: add_textbox(slide, x+3.35, 2.68, 0.35, 0.22, '→', 22, GREY, True)
quote(slide, 1.0, 4.55, 11.2, 1.0, 'La pregunta directiva no es solamente qué IA vamos a usar. Es si tenemos la información, los datos y el gobierno necesarios para que la IA cree valor real.', FUCHSIA)
add_textbox(slide, 1.15, 6.05, 10.9, 0.35, 'EXCELENCIA prepara la base confiable para IA y diseña agentes gobernados orientados a resultados.', 16, ANTH, True, align=PP_ALIGN.CENTER)
add_footer(slide, 12)

# 13
slide = blank(); set_bg(slide); add_title(slide, 'Beneficios esperables en ISO y PNC', 'Valor tangible e intangible')
card(slide, 0.85, 2.05, 2.75, 2.15, 'Tiempo experto', 'Menos tiempo en búsqueda documental, lectura repetitiva, comparación y armado de matrices.', TEAL, 15, 12.2)
card(slide, 3.85, 2.05, 2.75, 2.15, 'Consistencia', 'Mismo criterio entre filas, menor variabilidad y estructura homogénea de análisis.', ORANGE, 15, 12.2)
card(slide, 6.85, 2.05, 2.75, 2.15, 'Trazabilidad', 'Requisito o criterio → evidencia → fuente → juicio experto → acción.', FUCHSIA, 15, 12.2)
card(slide, 9.85, 2.05, 2.75, 2.15, 'Foco directivo', 'Más tiempo para validar, priorizar brechas, decidir y orientar planes de mejora.', YELLOW, 15, 12.2)
quote(slide, 0.95, 5.1, 11.45, 0.92, 'El beneficio no es reemplazar al experto: es devolverle tiempo y foco para decidir mejor.', TEAL)
add_textbox(slide, 1.05, 6.25, 11.2, 0.28, 'Aplicado a ISO y PNC: acelerar el ordenamiento de evidencia sin resignar criterio, trazabilidad ni validación humana.', 12.8, GREY, False, align=PP_ALIGN.CENTER)
add_footer(slide, 13)

# 14
slide = blank(); set_bg(slide)
diagonal_band(slide, 9.45, -0.15, 0.65, 7.95, YELLOW, 10)
diagonal_band(slide, 10.05, -0.15, 0.65, 7.95, ORANGE, 10)
diagonal_band(slide, 10.65, -0.15, 0.65, 7.95, FUCHSIA, 10)
diagonal_band(slide, 11.25, -0.15, 0.65, 7.95, TEAL, 10)
add_logo(slide)
add_textbox(slide, 0.82, 0.85, 4.5, 0.3, 'PREGUNTA FINAL', 11, TEAL, True)
add_textbox(slide, 0.82, 1.55, 8.0, 2.0, 'Si mañana incorporáramos IA en nuestros procesos críticos, ¿trabajaría sobre la mejor versión de nuestra organización o sobre nuestro desorden documental?', 30, ANTH, True, line_spacing=0.95)
quote(slide, 0.88, 4.65, 7.6, 0.82, 'Empezar por un diagnóstico: información confiable, datos íntegros, procesos claros y gobierno suficiente para que la IA cree valor real.', FUCHSIA)
add_textbox(slide, 0.95, 6.18, 7.5, 0.36, 'EXCELENCIA · Preparar la base confiable para IA y diseñar agentes gobernados orientados a resultados.', 13.5, GREY, False)
add_footer(slide, 14)

prs.save(OUT)
print(OUT)
